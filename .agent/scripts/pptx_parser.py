import os
import sys
from pptx import Presentation

def parse_pptx(file_path, output_path=None):
    if not os.path.exists(file_path):
        print(f"Error: File not found at {file_path}")
        return False
        
    try:
        prs = Presentation(file_path)
        md_lines = []
        
        for slide_idx, slide in enumerate(prs.slides):
            md_lines.append(f"\n---")
            md_lines.append(f"## Slide {slide_idx + 1}")
            
            # Extract notes if any
            notes_slide = slide.notes_slide
            notes_text = notes_slide.notes_text_frame.text.strip() if notes_slide else ""
            
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            md_lines.append(text)
                elif shape.has_table:
                    md_lines.append("")
                    for r_idx, row in enumerate(shape.table.rows):
                        row_text = [cell.text.replace('\n', ' ').strip() for cell in row.cells]
                        md_lines.append("| " + " | ".join(row_text) + " |")
                        if r_idx == 0:
                            md_lines.append("| " + " | ".join(['---'] * len(row.cells)) + " |")
                    md_lines.append("")
            
            if notes_text:
                md_lines.append(f"\n*Speaker Notes:* {notes_text}")
                
        content = "\n".join(md_lines)
        if output_path:
            with open(output_path, 'w', encoding='utf-8') as out_f:
                out_f.write(content)
            print(f"Successfully converted {file_path} to {output_path}")
        else:
            print(content)
        return True
    except Exception as e:
        print(f"Error parsing PPTX: {e}")
        return False

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python pptx_parser.py <input_pptx_path> [output_md_path]")
        sys.exit(1)
        
    in_file = sys.argv[1]
    out_file = sys.argv[2] if len(sys.argv) > 2 else None
    parse_pptx(in_file, out_file)
